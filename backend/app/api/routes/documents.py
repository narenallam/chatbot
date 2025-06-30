"""
Document management API routes for the Personal Assistant AI Chatbot
Handles document preview, conversion, and source tracking
"""

from fastapi import APIRouter, HTTPException, Response
from fastapi.responses import StreamingResponse
from typing import List, Dict, Any, Optional
import logging
import io
import os
from pathlib import Path
import mimetypes

# Document conversion libraries
try:
    from docx import Document as DocxDocument
    from docx2pdf import convert as docx_to_pdf

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import inch

    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    import markdown
    from bs4 import BeautifulSoup

    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    import openpyxl
    import xlrd

    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from PIL import Image, ImageDraw, ImageFont
    import pillow_heif

    # Register HEIF opener with Pillow
    pillow_heif.register_heif_opener()
    IMAGE_AVAILABLE = True
except ImportError:
    IMAGE_AVAILABLE = False

from app.services.document_service import document_service

logger = logging.getLogger(__name__)

router = APIRouter()


def text_to_pdf(text_content: str, title: str = "Document") -> bytes:
    """Convert plain text to PDF using ReportLab"""
    buffer = io.BytesIO()

    try:
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()

        # Create custom styles
        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Title"],
            fontSize=16,
            spaceAfter=30,
        )

        normal_style = ParagraphStyle(
            "CustomNormal",
            parent=styles["Normal"],
            fontSize=11,
            spaceAfter=12,
            leftIndent=0,
            rightIndent=0,
        )

        # Build document content
        story = []

        # Add title
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 12))

        # Add content, handling line breaks
        lines = text_content.split("\n")
        for line in lines:
            if line.strip():
                # Escape HTML characters and handle special characters
                # Also clean up problematic Unicode characters
                line = (
                    line.replace("&", "&amp;")
                    .replace("<", "&lt;")
                    .replace(">", "&gt;")
                    .replace(
                        "\u202f", " "
                    )  # Replace narrow no-break space with regular space
                    .replace(
                        "\u00a0", " "
                    )  # Replace non-breaking space with regular space
                    .replace("\u2019", "'")  # Replace right single quotation mark
                    .replace("\u2018", "'")  # Replace left single quotation mark
                    .replace("\u201c", '"')  # Replace left double quotation mark
                    .replace("\u201d", '"')  # Replace right double quotation mark
                    .replace("\u2013", "-")  # Replace en dash
                    .replace("\u2014", "-")  # Replace em dash
                )

                # Ensure the line can be encoded safely for ReportLab
                try:
                    line.encode("latin-1")
                except UnicodeEncodeError:
                    # If there are still problematic characters, encode/decode to clean them
                    line = line.encode("ascii", errors="ignore").decode("ascii")

                story.append(Paragraph(line, normal_style))
            else:
                story.append(Spacer(1, 6))

        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()
    except Exception as e:
        logger.error(f"Failed to convert text to PDF: {e}")
        # Fallback: create simple PDF with error message
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = [
            Paragraph(f"Error converting document: {title}", styles["Title"]),
            Paragraph(f"Error: {str(e)}", styles["Normal"]),
        ]
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()


def markdown_to_pdf(markdown_content: str, title: str = "Document") -> bytes:
    """Convert Markdown to PDF"""
    if not MARKDOWN_AVAILABLE:
        return text_to_pdf(markdown_content, title)

    try:
        # Convert markdown to HTML
        html = markdown.markdown(markdown_content)

        # Parse HTML and convert to plain text for PDF
        soup = BeautifulSoup(html, "html.parser")
        text_content = soup.get_text()

        return text_to_pdf(text_content, title)
    except Exception as e:
        logger.error(f"Failed to convert markdown to PDF: {e}")
        return text_to_pdf(markdown_content, title)


def html_to_pdf(html_content: str, title: str = "Document") -> bytes:
    """Convert HTML to PDF"""
    try:
        # Parse HTML and convert to plain text for PDF
        soup = BeautifulSoup(html_content, "html.parser")
        text_content = soup.get_text()

        return text_to_pdf(text_content, title)
    except Exception as e:
        logger.error(f"Failed to convert HTML to PDF: {e}")
        return text_to_pdf(html_content, title)


def pptx_to_pdf(pptx_content: bytes, title: str = "Presentation") -> bytes:
    """Convert PowerPoint to PDF"""
    if not PPTX_AVAILABLE:
        return text_to_pdf("PowerPoint conversion not available", title)

    try:
        pptx_file = io.BytesIO(pptx_content)
        presentation = Presentation(pptx_file)

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()

        title_style = ParagraphStyle(
            "PresentationTitle",
            parent=styles["Title"],
            fontSize=18,
            spaceAfter=30,
        )

        slide_title_style = ParagraphStyle(
            "SlideTitle",
            parent=styles["Heading1"],
            fontSize=14,
            spaceAfter=12,
            textColor="blue",
        )

        content_style = ParagraphStyle(
            "SlideContent",
            parent=styles["Normal"],
            fontSize=10,
            spaceAfter=8,
            leftIndent=20,
        )

        story = []
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 20))

        for i, slide in enumerate(presentation.slides):
            # Add slide number
            story.append(Paragraph(f"Slide {i + 1}", slide_title_style))

            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                for text in slide_text:
                    # Escape HTML characters
                    text = (
                        text.replace("&", "&amp;")
                        .replace("<", "&lt;")
                        .replace(">", "&gt;")
                    )
                    story.append(Paragraph(text, content_style))
            else:
                story.append(Paragraph("(No text content)", content_style))

            story.append(Spacer(1, 20))

        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()

    except Exception as e:
        logger.error(f"Failed to convert PowerPoint to PDF: {e}")
        return text_to_pdf(f"Error converting PowerPoint: {str(e)}", title)


def excel_to_pdf(excel_content: bytes, title: str = "Spreadsheet") -> bytes:
    """Convert Excel to PDF"""
    if not EXCEL_AVAILABLE:
        return text_to_pdf("Excel conversion not available", title)

    try:
        excel_file = io.BytesIO(excel_content)

        # Try to read as Excel file
        try:
            # Read all sheets
            sheets_dict = pd.read_excel(excel_file, sheet_name=None, engine="openpyxl")
        except:
            try:
                # Try with xlrd engine for older .xls files
                sheets_dict = pd.read_excel(excel_file, sheet_name=None, engine="xlrd")
            except:
                return text_to_pdf("Could not read Excel file", title)

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()

        title_style = ParagraphStyle(
            "SpreadsheetTitle",
            parent=styles["Title"],
            fontSize=16,
            spaceAfter=30,
        )

        sheet_title_style = ParagraphStyle(
            "SheetTitle",
            parent=styles["Heading2"],
            fontSize=12,
            spaceAfter=12,
            textColor="green",
        )

        cell_style = ParagraphStyle(
            "CellContent",
            parent=styles["Normal"],
            fontSize=8,
            spaceAfter=4,
        )

        story = []
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 20))

        for sheet_name, df in sheets_dict.items():
            story.append(Paragraph(f"Sheet: {sheet_name}", sheet_title_style))

            # Convert DataFrame to text representation
            if not df.empty:
                # Add column headers
                headers = " | ".join(str(col) for col in df.columns)
                story.append(Paragraph(f"Columns: {headers}", cell_style))
                story.append(Spacer(1, 6))

                # Add first few rows of data
                for idx, row in df.head(20).iterrows():  # Limit to first 20 rows
                    row_text = " | ".join(str(val) for val in row.values)
                    # Escape HTML and limit length
                    row_text = (
                        row_text.replace("&", "&amp;")
                        .replace("<", "&lt;")
                        .replace(">", "&gt;")
                    )
                    if len(row_text) > 100:
                        row_text = row_text[:97] + "..."
                    story.append(Paragraph(f"Row {idx + 1}: {row_text}", cell_style))

                if len(df) > 20:
                    story.append(
                        Paragraph(f"... and {len(df) - 20} more rows", cell_style)
                    )
            else:
                story.append(Paragraph("(Empty sheet)", cell_style))

            story.append(Spacer(1, 20))

        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()

    except Exception as e:
        logger.error(f"Failed to convert Excel to PDF: {e}")
        return text_to_pdf(f"Error converting Excel: {str(e)}", title)


def image_to_pdf(image_content: bytes, title: str = "Image") -> bytes:
    """Convert Image to PDF"""
    if not IMAGE_AVAILABLE:
        return text_to_pdf("Image conversion not available", title)

    try:
        image_file = io.BytesIO(image_content)
        image = Image.open(image_file)

        # Convert to RGB if necessary
        if image.mode != "RGB":
            image = image.convert("RGB")

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()

        title_style = ParagraphStyle(
            "ImageTitle",
            parent=styles["Title"],
            fontSize=16,
            spaceAfter=20,
        )

        story = []
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 20))

        # Add image information
        info_style = styles["Normal"]
        story.append(
            Paragraph(
                f"Image Size: {image.size[0]} x {image.size[1]} pixels", info_style
            )
        )
        story.append(Paragraph(f"Image Mode: {image.mode}", info_style))
        story.append(Spacer(1, 20))

        # Save image as JPEG for embedding
        img_buffer = io.BytesIO()
        image.save(img_buffer, format="JPEG", quality=85)
        img_buffer.seek(0)

        # Add image to PDF
        from reportlab.platypus import Image as ReportLabImage
        from reportlab.lib.units import inch

        # Calculate image size to fit on page
        page_width = A4[0] - 2 * inch  # Leave margins
        page_height = A4[1] - 4 * inch  # Leave space for title and margins

        img_width, img_height = image.size
        aspect_ratio = img_width / img_height

        if aspect_ratio > 1:  # Landscape
            display_width = min(page_width, 6 * inch)
            display_height = display_width / aspect_ratio
        else:  # Portrait
            display_height = min(page_height, 8 * inch)
            display_width = display_height * aspect_ratio

        # Create temporary file for the image
        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as temp_file:
            image.save(temp_file.name, format="JPEG", quality=85)
            story.append(
                ReportLabImage(
                    temp_file.name, width=display_width, height=display_height
                )
            )

        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()

    except Exception as e:
        logger.error(f"Failed to convert image to PDF: {e}")
        return text_to_pdf(f"Error converting image: {str(e)}", title)


@router.get("/preview/{document_id}")
async def preview_document(document_id: str):
    """
    Get document preview as PDF
    Converts non-PDF documents to PDF for consistent preview experience
    """
    try:
        # Get document info
        document_info = await document_service.get_document_info(document_id)
        if not document_info:
            raise HTTPException(status_code=404, detail="Document not found")

        # Get document metadata
        metadata = document_info.get("metadata", {})
        filename = metadata.get("filename", "unknown")
        file_path = metadata.get("file_path")
        content_type = metadata.get("content_type", "application/octet-stream")

        # If no file path, try to get content from chunks
        if not file_path or not os.path.exists(file_path):
            # Reconstruct content from chunks
            chunks = document_info.get("chunks", [])
            if not chunks:
                raise HTTPException(
                    status_code=404, detail="Document content not found"
                )

            content = "\n\n".join([chunk.get("content", "") for chunk in chunks])
            title = filename or f"Document {document_id}"

            # Convert to PDF
            pdf_content = text_to_pdf(content, title)

            return StreamingResponse(
                io.BytesIO(pdf_content),
                media_type="application/pdf",
                headers={"Content-Disposition": f'inline; filename="{title}.pdf"'},
            )

        # Handle file based on type
        file_extension = Path(filename).suffix.lower()

        if file_extension == ".pdf":
            # Return PDF directly
            with open(file_path, "rb") as f:
                pdf_content = f.read()

            return StreamingResponse(
                io.BytesIO(pdf_content),
                media_type="application/pdf",
                headers={"Content-Disposition": f'inline; filename="{filename}"'},
            )

        # Convert non-PDF files to PDF
        with open(file_path, "rb") as f:
            content = f.read()

        title = Path(filename).stem

        if file_extension == ".md":
            pdf_content = markdown_to_pdf(
                content.decode("utf-8", errors="ignore"), title
            )
        elif file_extension in [".html", ".htm"]:
            pdf_content = html_to_pdf(content.decode("utf-8", errors="ignore"), title)
        elif file_extension in [".txt", ".log"]:
            pdf_content = text_to_pdf(content.decode("utf-8", errors="ignore"), title)
        elif file_extension in [".ppt", ".pptx"]:
            pdf_content = pptx_to_pdf(content, title)
        elif file_extension in [".xls", ".xlsx"]:
            pdf_content = excel_to_pdf(content, title)
        elif file_extension in [
            ".png",
            ".jpg",
            ".jpeg",
            ".heic",
            ".bmp",
            ".gif",
            ".tiff",
        ]:
            pdf_content = image_to_pdf(content, title)
        else:
            # Default: treat as plain text
            pdf_content = text_to_pdf(content.decode("utf-8", errors="ignore"), title)

        return StreamingResponse(
            io.BytesIO(pdf_content),
            media_type="application/pdf",
            headers={"Content-Disposition": f'inline; filename="{title}.pdf"'},
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Preview document error: {e}")
        raise HTTPException(
            status_code=500, detail=f"Failed to preview document: {str(e)}"
        )


@router.get("/original/{document_id}")
async def get_original_document(document_id: str):
    """
    Get original document file (useful for images that don't need PDF conversion)
    """
    try:
        # Get document info
        document_info = await document_service.get_document_info(document_id)
        if not document_info:
            raise HTTPException(status_code=404, detail="Document not found")

        # Get document metadata
        metadata = document_info.get("metadata", {})
        filename = metadata.get("filename", "unknown")
        file_path = metadata.get("file_path")
        content_type = metadata.get("content_type", "application/octet-stream")

        # Check if file exists
        if not file_path or not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="Original file not found")

        # Read and return original file
        with open(file_path, "rb") as f:
            content = f.read()

        # Encode filename properly for HTTP headers
        safe_filename = filename.encode("ascii", "ignore").decode("ascii")
        if not safe_filename:
            safe_filename = "document"

        return StreamingResponse(
            io.BytesIO(content),
            media_type=content_type,
            headers={"Content-Disposition": f'inline; filename="{safe_filename}"'},
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Get original document error: {e}")
        raise HTTPException(
            status_code=500, detail=f"Failed to get original document: {str(e)}"
        )


@router.get("/download/{document_id}")
async def download_document(document_id: str):
    """
    Download document file (forces download instead of preview)
    """
    try:
        # Get document info
        document_info = await document_service.get_document_info(document_id)
        if not document_info:
            raise HTTPException(status_code=404, detail="Document not found")

        # Get document metadata
        metadata = document_info.get("metadata", {})
        filename = metadata.get("filename", "unknown")
        file_path = metadata.get("file_path")
        content_type = metadata.get("content_type", "application/octet-stream")

        # Check if file exists
        if not file_path or not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="File not found")

        # Read and return file with download headers
        with open(file_path, "rb") as f:
            content = f.read()

        # Encode filename properly for HTTP headers
        safe_filename = filename.encode("ascii", "ignore").decode("ascii")
        if not safe_filename:
            safe_filename = "document"

        return StreamingResponse(
            io.BytesIO(content),
            media_type=content_type,
            headers={"Content-Disposition": f'attachment; filename="{safe_filename}"'},
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Download document error: {e}")
        raise HTTPException(
            status_code=500, detail=f"Failed to download document: {str(e)}"
        )


@router.get("/download-pdf/{document_id}")
async def download_document_as_pdf(document_id: str):
    """
    Download document as PDF (forces download instead of preview)
    """
    try:
        # Get document info
        document_info = await document_service.get_document_info(document_id)
        if not document_info:
            raise HTTPException(status_code=404, detail="Document not found")

        # Get document metadata
        metadata = document_info.get("metadata", {})
        filename = metadata.get("filename", "unknown")
        file_path = metadata.get("file_path")

        # Generate PDF filename with safe encoding
        safe_filename_stem = filename.encode("ascii", "ignore").decode("ascii")
        if not safe_filename_stem:
            safe_filename_stem = "document"
        pdf_filename = f"{Path(safe_filename_stem).stem}.pdf"

        # If no file path, try to get content from chunks
        if not file_path or not os.path.exists(file_path):
            # Reconstruct content from chunks
            chunks = document_info.get("chunks", [])
            if not chunks:
                raise HTTPException(
                    status_code=404, detail="Document content not found"
                )

            content = "\n\n".join([chunk.get("content", "") for chunk in chunks])
            title = filename or f"Document {document_id}"

            # Convert to PDF
            pdf_content = text_to_pdf(content, title)

            return StreamingResponse(
                io.BytesIO(pdf_content),
                media_type="application/pdf",
                headers={
                    "Content-Disposition": f'attachment; filename="{pdf_filename}"'
                },
            )

        # Handle file based on type
        file_extension = Path(filename).suffix.lower()

        if file_extension == ".pdf":
            # Return PDF directly
            with open(file_path, "rb") as f:
                pdf_content = f.read()

            # Encode filename properly for HTTP headers
            safe_filename = filename.encode("ascii", "ignore").decode("ascii")
            if not safe_filename:
                safe_filename = "document.pdf"

            return StreamingResponse(
                io.BytesIO(pdf_content),
                media_type="application/pdf",
                headers={
                    "Content-Disposition": f'attachment; filename="{safe_filename}"'
                },
            )

        # Convert non-PDF files to PDF
        with open(file_path, "rb") as f:
            content = f.read()

        title = Path(filename).stem

        if file_extension == ".md":
            pdf_content = markdown_to_pdf(
                content.decode("utf-8", errors="ignore"), title
            )
        elif file_extension in [".html", ".htm"]:
            pdf_content = html_to_pdf(content.decode("utf-8", errors="ignore"), title)
        elif file_extension in [".txt", ".log"]:
            pdf_content = text_to_pdf(content.decode("utf-8", errors="ignore"), title)
        elif file_extension in [".ppt", ".pptx"]:
            pdf_content = pptx_to_pdf(content, title)
        elif file_extension in [".xls", ".xlsx"]:
            pdf_content = excel_to_pdf(content, title)
        elif file_extension in [
            ".png",
            ".jpg",
            ".jpeg",
            ".heic",
            ".bmp",
            ".gif",
            ".tiff",
        ]:
            pdf_content = image_to_pdf(content, title)
        else:
            # Default: treat as plain text
            pdf_content = text_to_pdf(content.decode("utf-8", errors="ignore"), title)

        return StreamingResponse(
            io.BytesIO(pdf_content),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{pdf_filename}"'},
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Download PDF error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to download PDF: {str(e)}")


@router.get("/sources/{conversation_id}")
async def get_conversation_sources(conversation_id: str):
    """
    Get document sources referenced in a conversation
    """
    try:
        # This would typically query the chat service for document sources
        # For now, return a placeholder response
        return {
            "conversation_id": conversation_id,
            "sources": [],
            "message": "Document source tracking not yet implemented",
        }
    except Exception as e:
        logger.error(f"Get conversation sources error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to get sources: {str(e)}")


@router.get("/search")
async def search_documents(query: str, limit: int = 10):
    """
    Search documents by content
    """
    try:
        # Use the document service to search
        results = await document_service.search_documents(query, limit=limit)

        return {"query": query, "results": results, "total": len(results)}
    except Exception as e:
        logger.error(f"Search documents error: {e}")
        raise HTTPException(
            status_code=500, detail=f"Failed to search documents: {str(e)}"
        )


@router.get("/")
async def list_documents(response: Response):
    """
    List all uploaded documents
    """
    try:
        # Add cache-busting headers
        response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"

        logger.info("=== DEBUG: /list endpoint called ===")
        raw_documents = await document_service.list_documents()
        logger.info(f"=== DEBUG: Got {len(raw_documents)} documents ===")

        # Map database fields to frontend-expected format
        documents = []
        for doc in raw_documents:
            mapped_doc = {
                "id": doc["id"],
                "name": doc["full_filename"],  # Map full_filename to name for frontend
                "size": doc["file_size"],
                "type": doc["content_type"],
                "uploadedAt": doc["uploaded_datetime"],
                "status": doc.get("status", "ready"),
                # Enhanced fields from the new document service
                "fullFileName": doc["full_filename"],
                "fileHash": doc["file_hash"],
                "newFileName": doc["new_filename"],
                "fileDataHash": doc.get("file_data_hash"),
                "contentType": doc["content_type"],
                "metadata": doc.get("metadata", {}),
            }
            documents.append(mapped_doc)

        return {"documents": documents, "total": len(documents)}
    except Exception as e:
        logger.error(f"List documents error: {e}")
        logger.error(f"=== DEBUG: Exception in /list endpoint: {str(e)} ===")
        raise HTTPException(
            status_code=500, detail=f"Failed to list documents: {str(e)}"
        )


@router.get("/{document_id}/info")
async def get_document_detailed_info(document_id: str):
    """
    Get detailed document information including preview URL
    """
    try:
        document_info = await document_service.get_document_info(document_id)
        if not document_info:
            raise HTTPException(status_code=404, detail="Document not found")

        # Add preview URL
        document_info["preview_url"] = f"/api/documents/preview/{document_id}"

        return document_info
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Get document info error: {e}")
        raise HTTPException(
            status_code=500, detail=f"Failed to get document info: {str(e)}"
        )


@router.delete("/all")
async def delete_all_documents():
    """
    Delete all documents (used by reset system)
    """
    try:
        documents = await document_service.list_documents()
        deleted_count = 0

        for doc in documents:
            success = await document_service.delete_document(doc["id"])
            if success:
                deleted_count += 1

        return {
            "message": f"Deleted {deleted_count} documents",
            "deleted_count": deleted_count,
            "total_documents": len(documents),
        }
    except Exception as e:
        logger.error(f"Delete all documents error: {e}")
        raise HTTPException(
            status_code=500, detail=f"Failed to delete documents: {str(e)}"
        )
