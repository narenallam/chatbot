"""
Enhanced Document Processing Service for the Personal Assistant AI Chatbot
Handles document upload, parsing, OCR, and storage for specific file types
"""

import asyncio
import aiofiles
import uuid
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging
from datetime import datetime
import warnings
import io
import sys
from contextlib import contextmanager
import hashlib
import json
import tempfile

# Document processing imports
import PyPDF2
from docx import Document as DocxDocument

# Additional document processing imports
try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    import openpyxl

    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

# Image and OCR processing imports
try:
    from PIL import Image

    IMAGE_AVAILABLE = True
except ImportError:
    IMAGE_AVAILABLE = False

# HEIF support
try:
    import pillow_heif

    pillow_heif.register_heif_opener()
    HEIF_AVAILABLE = True
except ImportError:
    HEIF_AVAILABLE = False

# OCR support
try:
    import pytesseract

    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# Enhanced PDF processing
try:
    import fitz  # PyMuPDF

    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

# LangChain imports for text processing
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

# Internal imports
from app.core.config import settings
from app.services.vector_service import vector_service
from app.services.database_service import database_service

# Import multiprocessing service
try:
    from app.services.multiprocessing_service import multiprocessing_service

    MULTIPROCESSING_AVAILABLE = True
except ImportError:
    MULTIPROCESSING_AVAILABLE = False
    multiprocessing_service = None

logger = logging.getLogger(__name__)

# Global function to broadcast logs (will be set by main.py)
broadcast_log_func = None


async def log_to_websocket(level: str, message: str, details: dict = None):
    """Send log message to WebSocket clients"""
    if broadcast_log_func:
        try:
            await broadcast_log_func(level, message, details)
        except Exception as e:
            logger.error(f"Failed to broadcast log: {e}")


@contextmanager
def suppress_pdf_warnings():
    """Context manager to suppress PyPDF2 warnings and capture them"""
    old_stderr = sys.stderr
    sys.stderr = io.StringIO()
    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        try:
            yield
        finally:
            captured = sys.stderr.getvalue()
            sys.stderr = old_stderr
            if captured:
                logger.debug(f"PDF parsing messages: {captured[:200]}...")


class DocumentService:
    """Enhanced document processing and storage service with OCR support"""

    def __init__(self):
        self.db_service = database_service
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
            length_function=len,
        )
        
        # Initialize table-aware text splitter for enhanced table handling
        self.table_aware_splitter = TableAwareTextSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
            table_chunk_size=settings.chunk_size * 2  # Larger chunks for tables
        )
        self._setup_storage_directories()
        self._log_capabilities()

    def _setup_storage_directories(self):
        """Setup storage directories"""
        directories = [
            Path("data/hashed_files"),
            Path("data/original_files"),  # Store original files with original names
            Path("data/metadata"),
            Path("data/logs"),
            Path("data/temp"),
        ]
        for directory in directories:
            directory.mkdir(parents=True, exist_ok=True)

    def _log_capabilities(self):
        """Log available processing capabilities"""
        capabilities = []
        if IMAGE_AVAILABLE:
            capabilities.append("Image processing (PIL)")
        if HEIF_AVAILABLE:
            capabilities.append("HEIF support")
        if OCR_AVAILABLE:
            capabilities.append("OCR (pytesseract)")
        if PYMUPDF_AVAILABLE:
            capabilities.append("Enhanced PDF (PyMuPDF)")
        if PPTX_AVAILABLE:
            capabilities.append("PowerPoint processing")
        if EXCEL_AVAILABLE:
            capabilities.append("Excel processing")

        logger.info(f"Document processing capabilities: {', '.join(capabilities)}")

    def _calculate_file_hash(self, file_content: bytes) -> str:
        """Calculate SHA256 hash of file content"""
        return hashlib.sha256(file_content).hexdigest()

    def _calculate_data_hash(self, file_content: bytes, file_size: int) -> str:
        """Calculate hash of file data + size for duplicate detection"""
        data = file_content + str(file_size).encode()
        return hashlib.sha256(data).hexdigest()

    def _generate_new_filename(self, file_hash: str, original_filename: str) -> str:
        """Generate new filename using hash and original extension"""
        extension = Path(original_filename).suffix.lower()
        return f"{file_hash}{extension}"

    def _get_content_type(self, filename: str) -> str:
        """Get MIME content type from filename"""
        extension = Path(filename).suffix.lower()
        content_types = {
            ".pdf": "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".heic": "image/heic",
            ".heif": "image/heif",
            ".md": "text/markdown",
            ".txt": "text/plain",
            ".c": "text/x-c",
            ".cpp": "text/x-c++",
            ".java": "text/x-java",
            ".cs": "text/x-csharp",
            ".js": "application/javascript",
            ".ts": "application/typescript",
            ".go": "text/x-go",
            ".rs": "text/x-rustsrc",
            ".php": "text/x-php",
            ".pl": "text/x-perl",
            ".rb": "text/x-ruby",
            ".py": "text/x-python",
            ".swift": "text/x-swift",
            ".kt": "text/x-kotlin",
            ".scala": "text/x-scala",
            ".sh": "application/x-sh",
            ".bat": "application/x-bat",
            ".ps1": "application/x-powershell",
            ".html": "text/html",
            ".css": "text/css",
            ".json": "application/json",
            ".xml": "application/xml",
            ".yaml": "application/x-yaml",
            ".yml": "application/x-yaml",
            ".toml": "application/x-toml",
        }
        return content_types.get(extension, "application/octet-stream")

    def _is_supported_file_type(self, filename: str) -> bool:
        """Check if file type is supported"""
        supported_extensions = {
            ".pdf",
            ".docx",
            ".pptx",
            ".xlsx",
            ".png",
            ".jpg",
            ".jpeg",
            ".heic",
            ".heif",
            ".md",
            ".txt",
            ".c",
            ".cpp",
            ".java",
            ".cs",
            ".js",
            ".ts",
            ".go",
            ".rs",
            ".php",
            ".pl",
            ".rb",
            ".py",
            ".swift",
            ".kt",
            ".scala",
            ".sh",
            ".bat",
            ".ps1",
            ".html",
            ".css",
            ".json",
            ".xml",
            ".yaml",
            ".yml",
            ".toml",
        }
        extension = Path(filename).suffix.lower()
        return extension in supported_extensions

    async def process_uploaded_file(
        self, file_content: bytes, filename: str, content_type: str
    ) -> Dict[str, Any]:
        """
        Process uploaded file with enhanced OCR capabilities

        Args:
            file_content: File content as bytes
            filename: Original filename
            content_type: MIME content type

        Returns:
            Processing result dictionary
        """
        start_time = datetime.now()

        try:
            # Validate file type
            if not self._is_supported_file_type(filename):
                return {
                    "status": "error",
                    "message": f"Unsupported file type: {Path(filename).suffix}",
                    "filename": filename,
                }

            # Calculate hashes
            file_hash = self._calculate_file_hash(file_content)
            file_size = len(file_content)
            file_data_hash = self._calculate_data_hash(file_content, file_size)

            # Check for duplicates
            existing_file = self.db_service.check_file_exists(file_data_hash)
            if existing_file:
                return {
                    "status": "success",
                    "message": "File already exists",
                    "filename": filename,
                    "file_id": existing_file["id"],
                    "duplicate": True,
                }

            # Generate new filename
            new_filename = self._generate_new_filename(file_hash, filename)

            # Save file to hashed_files directory
            file_path = Path("data/hashed_files") / new_filename
            async with aiofiles.open(file_path, "wb") as f:
                await f.write(file_content)

            # Also save original file with original name in original_files directory
            # Generate a unique original filename to handle duplicates
            original_file_path = Path("data/original_files") / filename
            counter = 1
            while original_file_path.exists():
                name_parts = Path(filename).stem, Path(filename).suffix
                original_file_path = (
                    Path("data/original_files")
                    / f"{name_parts[0]}_{counter}{name_parts[1]}"
                )
                counter += 1

            async with aiofiles.open(original_file_path, "wb") as f:
                await f.write(file_content)

            # Extract text content with OCR support
            text_content = await self._extract_text(
                file_content, filename, content_type
            )

            if not text_content.strip():
                return {
                    "status": "error",
                    "message": "No text content could be extracted from the file",
                    "filename": filename,
                }

            # Generate file ID first
            file_id = str(uuid.uuid4())

            # Create chunks for vector storage
            chunks = self._create_chunks(text_content, filename, file_hash, file_id)

            # Store in vector database
            chunk_ids = await vector_service.add_documents(chunks)

            # Save file info to database with the pre-generated file_id
            saved_file_id = self.db_service.save_file_info(
                file_id=file_id,
                full_filename=filename,
                file_hash=file_hash,
                new_filename=new_filename,
                file_size=file_size,
                file_data_hash=file_data_hash,
                content_type=content_type,
                metadata={
                    "chunk_count": len(chunks),
                    "chunk_ids": chunk_ids,
                    "processing_time": (datetime.now() - start_time).total_seconds(),
                    "ocr_used": self._was_ocr_used(filename, text_content),
                    "original_file_path": str(original_file_path),
                    "hashed_file_path": str(file_path),
                },
            )

            processing_time = (datetime.now() - start_time).total_seconds()

            await log_to_websocket(
                "info",
                f"✅ Successfully processed {filename}",
                {
                    "chunks": len(chunks),
                    "processing_time": processing_time,
                    "file_size_mb": file_size / 1024 / 1024,
                    "ocr_used": self._was_ocr_used(filename, text_content),
                },
            )

            return {
                "status": "success",
                "message": "File processed successfully",
                "filename": filename,
                "file_id": file_id,
                "chunk_count": len(chunks),
                "processing_time_seconds": processing_time,
                "file_size_mb": file_size / 1024 / 1024,
                "new_filename": new_filename,
                "ocr_used": self._was_ocr_used(filename, text_content),
            }

        except Exception as e:
            logger.error(f"Error processing file {filename}: {str(e)}")
            await log_to_websocket("error", f"❌ Error processing {filename}: {str(e)}")

            return {
                "status": "error",
                "message": f"Processing failed: {str(e)}",
                "filename": filename,
            }

    def _was_ocr_used(self, filename: str, text_content: str) -> bool:
        """Check if OCR was likely used based on filename and content"""
        extension = Path(filename).suffix.lower()
        is_image = extension in [".png", ".jpg", ".jpeg", ".heic", ".heif"]

        # If it's an image and has substantial text content, OCR was likely used
        if is_image and len(text_content.strip()) > 50:
            return True

        # Check for OCR indicators in text
        ocr_indicators = ["[OCR]", "scanned", "image text", "extracted from image"]
        return any(
            indicator.lower() in text_content.lower() for indicator in ocr_indicators
        )

    async def _extract_text(
        self, file_content: bytes, filename: str, content_type: str
    ) -> str:
        """
        Extract text from file based on type with OCR support
        """
        file_extension = Path(filename).suffix.lower()

        try:
            if file_extension == ".pdf":
                return await self._extract_from_pdf(file_content)
            elif file_extension == ".docx":
                return await self._extract_from_docx(file_content)
            elif file_extension == ".pptx":
                return await self._extract_from_pptx(file_content)
            elif file_extension == ".xlsx":
                return await self._extract_from_excel(file_content)
            elif file_extension in [".png", ".jpg", ".jpeg"]:
                return await self._extract_from_image(file_content, filename)
            elif file_extension in [".heic", ".heif"]:
                # Convert HEIC/HEIF to PNG using pillow-heif, then OCR
                if not HEIF_AVAILABLE:
                    raise ImportError(
                        "pillow-heif is not available for HEIC/HEIF conversion"
                    )
                from PIL import Image
                import io
                import pillow_heif

                heif_file = pillow_heif.read_heif(file_content)
                image = Image.frombytes(
                    heif_file.mode, heif_file.size, heif_file.data, "raw"
                )
                # Convert to PNG in memory
                png_bytes_io = io.BytesIO()
                image.save(png_bytes_io, format="PNG")
                png_bytes = png_bytes_io.getvalue()

                # Save the PNG file to hashed_files and original_files
                file_hash = self._calculate_file_hash(file_content)
                png_filename = f"{file_hash}.png"
                hashed_png_path = Path("data/hashed_files") / png_filename
                original_png_path = Path("data/original_files") / png_filename
                async with aiofiles.open(hashed_png_path, "wb") as f:
                    await f.write(png_bytes)
                async with aiofiles.open(original_png_path, "wb") as f:
                    await f.write(png_bytes)

                # Add PNG filename to metadata for preview
                if not hasattr(self, "_extra_metadata"):
                    self._extra_metadata = {}
                self._extra_metadata["converted_png"] = png_filename

                return await self._extract_from_image(png_bytes, png_filename)
            elif file_extension in [
                ".md",
                ".txt",
                ".c",
                ".cpp",
                ".java",
                ".cs",
                ".js",
                ".ts",
                ".go",
                ".rs",
                ".php",
                ".pl",
                ".rb",
                ".py",
                ".swift",
                ".kt",
                ".scala",
                ".sh",
                ".bat",
                ".ps1",
                ".html",
                ".css",
                ".json",
                ".xml",
                ".yaml",
                ".yml",
                ".toml",
            ]:
                # Treat as plain text/code/markdown
                try:
                    return file_content.decode("utf-8")
                except UnicodeDecodeError:
                    return file_content.decode("latin-1")
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")

        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {str(e)}")
            raise

    async def _extract_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file with OCR support and multiprocessing for large files"""
        try:
            # Check if multiprocessing should be used
            use_multiprocessing = (
                MULTIPROCESSING_AVAILABLE
                and PYMUPDF_AVAILABLE
                and settings.enable_parallel_processing
            )

            if use_multiprocessing:
                # Estimate PDF size to decide on multiprocessing
                file_size_mb = len(file_content) / (1024 * 1024)

                if file_size_mb > 5:  # Use multiprocessing for PDFs > 5MB
                    try:
                        logger.info(
                            f"Using multiprocessing for PDF ({file_size_mb:.1f}MB)"
                        )
                        await log_to_websocket(
                            "info",
                            f"🚀 Using parallel PDF processing ({file_size_mb:.1f}MB)",
                            {
                                "file_size_mb": file_size_mb,
                                "processing_type": "parallel",
                            },
                        )

                        extracted_text = (
                            await multiprocessing_service.process_pdf_parallel(
                                file_content, "uploaded_pdf"
                            )
                        )

                        if extracted_text.strip():
                            return extracted_text
                        else:
                            logger.warning(
                                "Parallel PDF processing returned empty text, falling back to sequential"
                            )
                    except Exception as e:
                        logger.warning(
                            f"Parallel PDF processing failed: {e}, falling back to sequential"
                        )

            # Fallback to sequential processing
            text_content = []

            # Try PyMuPDF first for better handling
            if PYMUPDF_AVAILABLE:
                try:
                    doc = fitz.open(stream=file_content, filetype="pdf")
                    total_pages = len(doc)

                    await log_to_websocket(
                        "info",
                        f"🔍 Processing {total_pages} pages with PyMuPDF (sequential)",
                        {"total_pages": total_pages, "processing_type": "sequential"},
                    )

                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)

                        # Extract text content and table data
                        page_content = await self._extract_page_content_with_tables(page, page_num)

                        # If no content found and OCR is available, try OCR
                        if not page_content.strip() and OCR_AVAILABLE:
                            await log_to_websocket(
                                "info", f"🔍 Using OCR for page {page_num + 1}"
                            )

                            # Convert page to image
                            pix = page.get_pixmap()
                            img_data = pix.tobytes("png")

                            # Perform OCR on the image
                            img = Image.open(io.BytesIO(img_data))
                            page_text = pytesseract.image_to_string(img)

                            if page_text.strip():
                                page_content = f"[OCR Page {page_num + 1}]\n{page_text}"

                        if page_content.strip():
                            text_content.append(f"[Page {page_num + 1}]\n{page_content}")

                    doc.close()

                    if text_content:
                        return "\n\n".join(text_content)

                except Exception as e:
                    logger.warning(
                        f"PyMuPDF extraction failed: {e}, falling back to PyPDF2"
                    )

            # Fallback to PyPDF2
            with suppress_pdf_warnings():
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))

                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append(f"[Page {page_num + 1}]\n{page_text}")
                    except Exception as e:
                        logger.warning(
                            f"Error extracting text from PDF page {page_num + 1}: {e}"
                        )
                        continue

                return "\n\n".join(text_content)

        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            raise

    async def _extract_page_content_with_tables(self, page, page_num: int) -> str:
        """
        Extract text and table content from a PDF page with proper structure preservation
        
        Args:
            page: PyMuPDF page object
            page_num: Page number for logging
            
        Returns:
            Formatted text content with structured tables
        """
        try:
            content_parts = []
            
            # Extract regular text content
            page_text = page.get_text()
            
            # Extract tables with structure preservation
            try:
                tables = page.find_tables()
                
                if tables:
                    await log_to_websocket(
                        "info", 
                        f"📊 Found {len(tables)} table(s) on page {page_num + 1}",
                        {"page": page_num + 1, "table_count": len(tables)}
                    )
                    
                    # Process each table
                    for table_idx, table in enumerate(tables):
                        try:
                            # Extract table data
                            table_data = table.extract()
                            
                            if table_data and len(table_data) > 0:
                                # Format table with proper structure
                                formatted_table = self._format_table_structure(
                                    table_data, page_num + 1, table_idx + 1
                                )
                                content_parts.append(formatted_table)
                                
                                # Remove table area from regular text to avoid duplication
                                bbox = table.bbox
                                # This is a simplified approach - in production you'd want more sophisticated text removal
                                
                        except Exception as e:
                            logger.warning(f"Error extracting table {table_idx + 1} on page {page_num + 1}: {e}")
                            continue
                            
            except Exception as e:
                logger.warning(f"Error finding tables on page {page_num + 1}: {e}")
            
            # Add remaining text content (non-table text)
            if page_text.strip():
                content_parts.append(page_text)
            
            return "\n\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"Error extracting page content with tables: {e}")
            # Fallback to regular text extraction
            return page.get_text()

    def _format_table_structure(self, table_data: list, page_num: int, table_num: int) -> str:
        """
        Format extracted table data with proper structure preservation
        
        Args:
            table_data: List of lists containing table cell data
            page_num: Page number
            table_num: Table number on the page
            
        Returns:
            Formatted table string with structure
        """
        try:
            if not table_data or len(table_data) == 0:
                return ""
            
            formatted_parts = []
            formatted_parts.append(f"=== TABLE {table_num} (Page {page_num}) ===")
            
            # Identify header row (first non-empty row)
            header_row = None
            data_rows = []
            
            for row_idx, row in enumerate(table_data):
                # Clean row data
                cleaned_row = [str(cell).strip() if cell is not None else "" for cell in row]
                
                # Skip completely empty rows
                if not any(cleaned_row):
                    continue
                
                if header_row is None:
                    header_row = cleaned_row
                    formatted_parts.append(f"HEADERS: {' | '.join(cleaned_row)}")
                    formatted_parts.append("---")
                else:
                    data_rows.append(cleaned_row)
            
            # Format data rows
            for row_idx, row in enumerate(data_rows):
                row_text = " | ".join(row)
                formatted_parts.append(f"ROW {row_idx + 1}: {row_text}")
            
            # Add table summary
            formatted_parts.append(f"=== END TABLE {table_num} ({len(data_rows)} data rows) ===")
            
            return "\n".join(formatted_parts)
            
        except Exception as e:
            logger.error(f"Error formatting table structure: {e}")
            # Fallback to simple format
            try:
                simple_table = []
                for row in table_data:
                    if row:
                        cleaned_row = [str(cell).strip() if cell is not None else "" for cell in row]
                        if any(cleaned_row):
                            simple_table.append(" | ".join(cleaned_row))
                return f"=== TABLE {table_num} (Page {page_num}) ===\n" + "\n".join(simple_table)
            except:
                return f"=== TABLE {table_num} (Page {page_num}) === [Error formatting table]"

    def _format_docx_table_structure(self, table, table_num: int) -> str:
        """
        Format DOCX table with proper structure preservation
        
        Args:
            table: python-docx table object
            table_num: Table number in the document
            
        Returns:
            Formatted table string with structure
        """
        try:
            if not table.rows:
                return ""
            
            formatted_parts = []
            formatted_parts.append(f"=== DOCX TABLE {table_num} ===")
            
            # Extract all rows first
            all_rows = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_data.append(cell_text)
                all_rows.append(row_data)
            
            if not all_rows:
                return ""
            
            # Identify header row (first non-empty row)
            header_row = None
            data_rows = []
            
            for row_idx, row in enumerate(all_rows):
                # Skip completely empty rows
                if not any(cell.strip() for cell in row):
                    continue
                
                if header_row is None:
                    header_row = row
                    formatted_parts.append(f"HEADERS: {' | '.join(row)}")
                    formatted_parts.append("---")
                else:
                    data_rows.append(row)
            
            # Format data rows
            for row_idx, row in enumerate(data_rows):
                row_text = " | ".join(row)
                if row_text.strip():
                    formatted_parts.append(f"ROW {row_idx + 1}: {row_text}")
            
            # Add table summary
            formatted_parts.append(f"=== END DOCX TABLE {table_num} ({len(data_rows)} data rows) ===")
            
            return "\n".join(formatted_parts)
            
        except Exception as e:
            logger.error(f"Error formatting DOCX table structure: {e}")
            # Fallback to simple format
            try:
                simple_rows = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        simple_rows.append(" | ".join(row_text))
                
                if simple_rows:
                    return f"=== DOCX TABLE {table_num} ===\n" + "\n".join(simple_rows)
                return ""
            except:
                return f"=== DOCX TABLE {table_num} === [Error formatting table]"

    async def _extract_from_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX file including embedded images"""
        try:
            doc = DocxDocument(io.BytesIO(file_content))
            text_content = []

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)

            # Extract text from tables with enhanced structure
            for table_idx, table in enumerate(doc.tables):
                formatted_table = self._format_docx_table_structure(table, table_idx + 1)
                if formatted_table:
                    text_content.append(formatted_table)

            # Extract text from embedded images if OCR is available
            if OCR_AVAILABLE:
                image_text = await self._extract_from_docx_images(doc)
                if image_text:
                    text_content.append("\n=== Embedded Images ===\n")
                    text_content.append(image_text)

            return "\n".join(text_content)

        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            raise

    async def _extract_from_docx_images(self, doc) -> str:
        """Extract text from images embedded in DOCX using OCR"""
        try:
            image_texts = []

            # Extract images from the document
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        img = Image.open(io.BytesIO(image_data))

                        # Perform OCR on the image
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            image_texts.append(
                                f"[Embedded Image OCR]\n{ocr_text.strip()}"
                            )

                    except Exception as e:
                        logger.warning(f"Error processing embedded image: {e}")
                        continue

            return "\n\n".join(image_texts)

        except Exception as e:
            logger.warning(f"Error extracting from DOCX images: {e}")
            return ""

    async def _extract_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PPTX file including embedded images"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not available")

        try:
            prs = Presentation(io.BytesIO(file_content))
            text_content = []

            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"[Slide {slide_num + 1}]"]

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                # Extract text from embedded images if OCR is available
                if OCR_AVAILABLE:
                    image_text = await self._extract_from_pptx_images(slide)
                    if image_text:
                        slide_text.append(
                            f"[Slide {slide_num + 1} Images]\n{image_text}"
                        )

                if len(slide_text) > 1:  # More than just the slide number
                    text_content.append("\n".join(slide_text))

            return "\n\n".join(text_content)

        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            raise

    async def _extract_from_pptx_images(self, slide) -> str:
        """Extract text from images embedded in PPTX slide using OCR"""
        try:
            image_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        image_data = shape.image.blob
                        img = Image.open(io.BytesIO(image_data))

                        # Perform OCR on the image
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            image_texts.append(f"[Slide Image OCR]\n{ocr_text.strip()}")

                    except Exception as e:
                        logger.warning(f"Error processing slide image: {e}")
                        continue

            return "\n\n".join(image_texts)

        except Exception as e:
            logger.warning(f"Error extracting from PPTX images: {e}")
            return ""

    async def _extract_from_excel(self, file_content: bytes) -> str:
        """Extract text from Excel file with enhanced table structure processing"""
        if not EXCEL_AVAILABLE:
            raise ImportError("pandas/openpyxl is not available")

        try:
            # Try to read with pandas
            df_dict = pd.read_excel(io.BytesIO(file_content), sheet_name=None)
            text_content = []

            for sheet_name, df in df_dict.items():
                if not df.empty:
                    formatted_sheet = await self._format_excel_sheet_structure(df, sheet_name)
                    if formatted_sheet:
                        text_content.append(formatted_sheet)

            return "\n\n".join(text_content)

        except Exception as e:
            logger.error(f"Error extracting text from Excel: {str(e)}")
            raise

    async def _format_excel_sheet_structure(self, df, sheet_name: str) -> str:
        """
        Format Excel sheet with proper table structure and intelligent row handling
        
        Args:
            df: pandas DataFrame
            sheet_name: Name of the Excel sheet
            
        Returns:
            Formatted sheet content with structure preservation
        """
        try:
            if df.empty:
                return ""
            
            formatted_parts = []
            formatted_parts.append(f"=== EXCEL SHEET: {sheet_name} ===")
            
            # Add sheet statistics
            total_rows, total_cols = df.shape
            await log_to_websocket(
                "info", 
                f"📊 Processing Excel sheet '{sheet_name}': {total_rows} rows, {total_cols} columns",
                {"sheet": sheet_name, "rows": total_rows, "columns": total_cols}
            )
            
            # Add column headers
            if not df.columns.empty:
                headers = [str(col) for col in df.columns]
                formatted_parts.append(f"HEADERS: {' | '.join(headers)}")
                formatted_parts.append("---")
            
            # Process all rows (not just first 10)
            # For very large sheets, we'll process intelligently
            if total_rows <= 1000:
                # Process all rows for smaller sheets
                for idx, row in df.iterrows():
                    row_values = []
                    for val in row.values:
                        if pd.notna(val):
                            # Handle different data types properly
                            if isinstance(val, (int, float)):
                                if pd.isna(val):
                                    row_values.append("")
                                else:
                                    row_values.append(str(val))
                            else:
                                row_values.append(str(val).strip())
                        else:
                            row_values.append("")
                    
                    # Only include rows with actual data
                    if any(cell.strip() for cell in row_values):
                        row_text = " | ".join(row_values)
                        formatted_parts.append(f"ROW {idx + 1}: {row_text}")
            else:
                # For large sheets, use intelligent sampling
                await log_to_websocket(
                    "info", 
                    f"📈 Large sheet detected ({total_rows} rows). Using intelligent sampling strategy."
                )
                
                # Process first 50 rows
                for idx, row in df.head(50).iterrows():
                    row_values = []
                    for val in row.values:
                        if pd.notna(val):
                            row_values.append(str(val).strip())
                        else:
                            row_values.append("")
                    
                    if any(cell.strip() for cell in row_values):
                        row_text = " | ".join(row_values)
                        formatted_parts.append(f"ROW {idx + 1}: {row_text}")
                
                # Add middle sample (around middle of data)
                if total_rows > 100:
                    middle_start = total_rows // 2 - 25
                    middle_end = total_rows // 2 + 25
                    formatted_parts.append(f"\n--- MIDDLE SECTION (Rows {middle_start}-{middle_end}) ---")
                    
                    for idx, row in df.iloc[middle_start:middle_end].iterrows():
                        row_values = []
                        for val in row.values:
                            if pd.notna(val):
                                row_values.append(str(val).strip())
                            else:
                                row_values.append("")
                        
                        if any(cell.strip() for cell in row_values):
                            row_text = " | ".join(row_values)
                            formatted_parts.append(f"ROW {idx + 1}: {row_text}")
                
                # Process last 50 rows
                formatted_parts.append(f"\n--- LAST SECTION (Final 50 rows) ---")
                for idx, row in df.tail(50).iterrows():
                    row_values = []
                    for val in row.values:
                        if pd.notna(val):
                            row_values.append(str(val).strip())
                        else:
                            row_values.append("")
                    
                    if any(cell.strip() for cell in row_values):
                        row_text = " | ".join(row_values)
                        formatted_parts.append(f"ROW {idx + 1}: {row_text}")
            
            # Add summary statistics
            processed_rows = len([line for line in formatted_parts if line.startswith("ROW")])
            formatted_parts.append(f"=== END EXCEL SHEET: {sheet_name} ({processed_rows} data rows processed) ===")
            
            return "\n".join(formatted_parts)
            
        except Exception as e:
            logger.error(f"Error formatting Excel sheet structure: {e}")
            # Fallback to simple format
            try:
                simple_parts = [f"=== EXCEL SHEET: {sheet_name} ==="]
                if not df.columns.empty:
                    simple_parts.append("Headers: " + " | ".join(str(col) for col in df.columns))
                
                # Process first 20 rows as fallback
                for idx, row in df.head(20).iterrows():
                    row_text = " | ".join(str(val) for val in row.values if pd.notna(val))
                    if row_text.strip():
                        simple_parts.append(f"Row {idx + 1}: {row_text}")
                
                return "\n".join(simple_parts)
            except:
                return f"=== EXCEL SHEET: {sheet_name} === [Error formatting sheet]"

    async def _extract_from_image(self, file_content: bytes, filename: str) -> str:
        """Extract text from image file using enhanced OCR with table detection"""
        if not IMAGE_AVAILABLE:
            raise ImportError("Pillow is not available")

        try:
            # Open image (supports HEIC if pillow-heif is available)
            image = Image.open(io.BytesIO(file_content))

            # Convert to RGB if necessary for OCR
            if image.mode not in ["RGB", "L"]:
                image = image.convert("RGB")

            # Extract basic image metadata
            metadata = [
                f"[Image: {filename}]",
                f"Format: {image.format}",
                f"Size: {image.size[0]}x{image.size[1]} pixels",
                f"Mode: {image.mode}",
            ]

            # Enhanced OCR processing with table detection
            extracted_content = await self._extract_image_content_with_tables(image, filename)
            
            # Try to extract EXIF data if available
            if hasattr(image, "_getexif") and image._getexif():
                exif = image._getexif()
                if exif:
                    metadata.append("EXIF data available")

            # Combine metadata with extracted content
            if extracted_content:
                return "\n".join(metadata) + "\n\n" + extracted_content
            else:
                return "\n".join(metadata) + "\n=== No text content extracted ==="

        except Exception as e:
            logger.error(f"Error extracting metadata from image: {str(e)}")
            raise

    async def _extract_image_content_with_tables(self, image, filename: str) -> str:
        """
        Extract content from image with enhanced table detection and OCR
        
        Args:
            image: PIL Image object
            filename: Image filename for logging
            
        Returns:
            Extracted text content with structured table formatting
        """
        try:
            if not OCR_AVAILABLE:
                return "=== OCR not available ==="
            
            content_parts = []
            
            # Enhanced OCR approach for table detection
            await log_to_websocket(
                "info", 
                f"🔍 Performing enhanced OCR with table detection on {filename}"
            )
            
            # Step 1: Basic OCR extraction
            try:
                basic_text = pytesseract.image_to_string(image)
                if basic_text.strip():
                    await log_to_websocket(
                        "info", 
                        f"📝 Basic OCR extracted {len(basic_text)} characters from {filename}"
                    )
            except Exception as e:
                logger.warning(f"Basic OCR failed for {filename}: {e}")
                basic_text = ""
            
            # Step 2: Table-specific OCR using different PSM modes
            table_texts = []
            
            # PSM 6: Uniform block of text (good for tables)
            try:
                table_config = '--psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz .,|-_()[]{}$%'
                table_text = pytesseract.image_to_string(image, config=table_config)
                if table_text.strip() and len(table_text.strip()) > 20:
                    table_texts.append(("PSM 6 (Table mode)", table_text.strip()))
                    await log_to_websocket(
                        "info", 
                        f"📊 Table-mode OCR extracted {len(table_text)} characters"
                    )
            except Exception as e:
                logger.warning(f"Table-mode OCR failed: {e}")
            
            # Step 3: Detect potential table structure in extracted text
            structured_content = await self._structure_image_table_content(
                basic_text, table_texts, filename
            )
            
            if structured_content:
                content_parts.append(structured_content)
            elif basic_text.strip():
                content_parts.append(f"=== OCR Extracted Text ===\n{basic_text.strip()}")
            else:
                content_parts.append("=== OCR: No text detected ===")
            
            return "\n\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"Error in enhanced image OCR: {e}")
            # Fallback to basic OCR
            try:
                basic_text = pytesseract.image_to_string(image)
                if basic_text.strip():
                    return f"=== OCR Extracted Text (Fallback) ===\n{basic_text.strip()}"
                else:
                    return "=== OCR: No text detected (Fallback) ==="
            except:
                return f"=== OCR failed: {str(e)} ==="

    async def _structure_image_table_content(self, basic_text: str, table_texts: List[tuple], filename: str) -> str:
        """
        Analyze and structure potential table content from image OCR
        
        Args:
            basic_text: Basic OCR extracted text
            table_texts: List of (method, text) tuples from different OCR modes
            filename: Image filename for logging
            
        Returns:
            Structured table content or None if no table detected
        """
        try:
            import re
            
            # Heuristics to detect if the image contains a table
            table_indicators = [
                r'\|',  # Pipe separators
                r'\s{3,}',  # Multiple spaces (column separation)
                r'\d+\s+\d+',  # Numbers with spaces (tabular data)
                r'[A-Za-z]+\s{2,}[A-Za-z]+',  # Words with multiple spaces
                r'Total|Sum|Amount|Revenue|Sales|Year|Month|Date|Name|Price|Quantity',  # Table keywords
            ]
            
            # Check all extracted texts for table indicators
            all_texts = [basic_text] + [text for _, text in table_texts]
            table_likelihood = 0
            
            for text in all_texts:
                if not text:
                    continue
                    
                for pattern in table_indicators:
                    matches = len(re.findall(pattern, text, re.IGNORECASE))
                    table_likelihood += matches
            
            await log_to_websocket(
                "info", 
                f"📊 Table likelihood score for {filename}: {table_likelihood}"
            )
            
            # If high likelihood of table content, structure it
            if table_likelihood >= 5:  # Threshold for table detection
                return await self._format_image_table_structure(all_texts, filename)
            
            return None
            
        except Exception as e:
            logger.error(f"Error structuring image table content: {e}")
            return None

    async def _format_image_table_structure(self, texts: List[str], filename: str) -> str:
        """
        Format detected table content from image OCR with proper structure
        
        Args:
            texts: List of OCR extracted texts
            filename: Image filename
            
        Returns:
            Formatted table structure
        """
        try:
            import re
            
            # Choose the best text (longest with most structured content)
            best_text = ""
            best_score = 0
            
            for text in texts:
                if not text:
                    continue
                    
                # Score based on length and structure indicators
                score = len(text)
                score += len(re.findall(r'\|', text)) * 10  # Pipe separators
                score += len(re.findall(r'\s{3,}', text)) * 5  # Multiple spaces
                score += len(re.findall(r'\d+', text)) * 2  # Numbers
                
                if score > best_score:
                    best_score = score
                    best_text = text
            
            if not best_text:
                return None
            
            formatted_parts = []
            formatted_parts.append(f"=== IMAGE TABLE: {filename} ===")
            
            # Split into lines and process
            lines = best_text.split('\n')
            processed_lines = []
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Try to detect column separators and format as table rows
                if '|' in line:
                    # Already has pipe separators
                    cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                    if len(cells) >= 2:
                        processed_lines.append(" | ".join(cells))
                elif re.search(r'\s{3,}', line):
                    # Multiple spaces indicate columns
                    cells = re.split(r'\s{3,}', line)
                    cells = [cell.strip() for cell in cells if cell.strip()]
                    if len(cells) >= 2:
                        processed_lines.append(" | ".join(cells))
                elif line:
                    # Single line of text
                    processed_lines.append(line)
            
            # Identify header and data rows
            if processed_lines:
                # First line with multiple columns is likely header
                header_found = False
                for i, line in enumerate(processed_lines):
                    if '|' in line and not header_found:
                        formatted_parts.append(f"HEADERS: {line}")
                        formatted_parts.append("---")
                        header_found = True
                    elif '|' in line:
                        formatted_parts.append(f"ROW {i}: {line}")
                    else:
                        formatted_parts.append(line)
            
            formatted_parts.append(f"=== END IMAGE TABLE: {filename} ({len(processed_lines)} rows processed) ===")
            
            await log_to_websocket(
                "info", 
                f"📊 Structured table from {filename}: {len(processed_lines)} rows"
            )
            
            return "\n".join(formatted_parts)
            
        except Exception as e:
            logger.error(f"Error formatting image table structure: {e}")
            return None

    def _create_chunks(
        self, text_content: str, filename: str, file_hash: str, document_id: str
    ) -> List[Document]:
        """
        Create text chunks for vector storage

        Args:
            text_content: Full text content
            filename: Original filename
            file_hash: File hash for reference

        Returns:
            List of Document chunks
        """
        try:
            # Use table-aware chunking for better table handling
            chunks = self.table_aware_splitter.split_text(text_content)

            documents = []
            for i, chunk in enumerate(chunks):
                if chunk.strip():
                    doc = Document(
                        page_content=chunk,
                        metadata={
                            "source": filename,
                            "file_hash": file_hash,
                            "document_id": document_id,
                            "chunk_index": i,
                            "total_chunks": len(chunks),
                            "chunk_id": f"{file_hash}_chunk_{i}",
                        },
                    )
                    documents.append(doc)

            return documents

        except Exception as e:
            logger.error(f"Error creating chunks for {filename}: {str(e)}")
            raise

    async def get_document_info(self, file_id: str) -> Optional[Dict[str, Any]]:
        """Get document information by file ID"""
        try:
            files = self.db_service.get_files()
            for file in files:
                if file["id"] == file_id:
                    # Add file path information for preview
                    file_info = file.copy()

                    # Determine the file path based on storage structure
                    original_path = Path("data/original_files") / file["full_filename"]
                    hashed_path = Path("data/hashed_files") / file["new_filename"]

                    # Prefer original file if it exists, otherwise use hashed file
                    if original_path.exists():
                        file_info["file_path"] = str(original_path)
                    elif hashed_path.exists():
                        file_info["file_path"] = str(hashed_path)
                    else:
                        logger.warning(
                            f"File not found for document {file_id}: {file['full_filename']}"
                        )
                        file_info["file_path"] = None

                    # Add metadata for preview endpoint compatibility
                    if "metadata" not in file_info:
                        file_info["metadata"] = {}

                    file_info["metadata"]["filename"] = file["full_filename"]
                    file_info["metadata"]["file_path"] = file_info["file_path"]
                    file_info["metadata"]["content_type"] = file["content_type"]

                    return file_info
            return None
        except Exception as e:
            logger.error(f"Error getting document info: {str(e)}")
            return None

    async def list_documents(self) -> List[Dict[str, Any]]:
        """List all uploaded documents"""
        try:
            return self.db_service.get_files()
        except Exception as e:
            logger.error(f"Error listing documents: {str(e)}")
            return []

    async def delete_document(self, file_id: str) -> bool:
        """Delete a document"""
        try:
            # Get file info
            file_info = await self.get_document_info(file_id)
            if not file_info:
                return False

            # Delete from vector database
            if "chunk_ids" in file_info.get("metadata", {}):
                chunk_ids = file_info["metadata"]["chunk_ids"]
                await vector_service.delete_documents(chunk_ids)

            # Delete file from storage
            file_path = Path("data/hashed_files") / file_info["new_filename"]
            if file_path.exists():
                file_path.unlink()

            # Delete from database
            return self.db_service.delete_file(file_id)

        except Exception as e:
            logger.error(f"Error deleting document: {str(e)}")
            return False

    async def search_documents(
        self, query: str, limit: int = 10
    ) -> List[Dict[str, Any]]:
        """Search documents using vector similarity"""
        try:
            results = await vector_service.search(query, limit)
            return results
        except Exception as e:
            logger.error(f"Error searching documents: {str(e)}")
            return []


class TableAwareTextSplitter:
    """
    Enhanced text splitter that preserves table structure and handles structured content intelligently
    """
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200, table_chunk_size: int = 2000):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.table_chunk_size = table_chunk_size
        
        # Create fallback standard splitter
        self.standard_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
        )
    
    def split_text(self, text: str) -> List[str]:
        """
        Split text with table structure awareness
        
        Args:
            text: Input text containing potentially structured table data
            
        Returns:
            List of text chunks with preserved table structure
        """
        try:
            # Identify table sections
            table_sections = self._identify_table_sections(text)
            
            if not table_sections:
                # No tables found, use standard chunking
                return self.standard_splitter.split_text(text)
            
            # Process text with table awareness
            return self._split_with_table_awareness(text, table_sections)
            
        except Exception as e:
            logger.error(f"Error in table-aware splitting, falling back to standard: {e}")
            return self.standard_splitter.split_text(text)
    
    def _identify_table_sections(self, text: str) -> List[Dict[str, Any]]:
        """
        Identify table sections in the text using markers
        
        Args:
            text: Input text
            
        Returns:
            List of table section information with start/end positions
        """
        try:
            import re
            
            table_sections = []
            
            # Look for table markers we added in our formatting
            table_patterns = [
                r'=== TABLE \d+ \(Page \d+\) ===',
                r'=== DOCX TABLE \d+ ===',
                r'=== EXCEL SHEET: .+ ===',
            ]
            
            for pattern in table_patterns:
                matches = list(re.finditer(pattern, text, re.MULTILINE))
                
                for match in matches:
                    # Find the end of this table
                    start_pos = match.start()
                    
                    # Look for end marker
                    end_patterns = [
                        r'=== END TABLE \d+ \(\d+ data rows\) ===',
                        r'=== END DOCX TABLE \d+ \(\d+ data rows\) ===',
                        r'=== END EXCEL SHEET: .+ \(\d+ data rows processed\) ===',
                    ]
                    
                    end_pos = len(text)  # Default to end of text
                    
                    for end_pattern in end_patterns:
                        end_match = re.search(end_pattern, text[start_pos:])
                        if end_match:
                            end_pos = start_pos + end_match.end()
                            break
                    
                    table_sections.append({
                        'start': start_pos,
                        'end': end_pos,
                        'type': 'table',
                        'content': text[start_pos:end_pos]
                    })
            
            # Sort by start position
            table_sections.sort(key=lambda x: x['start'])
            
            return table_sections
            
        except Exception as e:
            logger.error(f"Error identifying table sections: {e}")
            return []
    
    def _split_with_table_awareness(self, text: str, table_sections: List[Dict[str, Any]]) -> List[str]:
        """
        Split text while preserving table structure
        
        Args:
            text: Input text
            table_sections: List of identified table sections
            
        Returns:
            List of chunks with preserved table structure
        """
        try:
            chunks = []
            current_pos = 0
            
            for table_section in table_sections:
                # Process text before this table with standard chunking
                before_table = text[current_pos:table_section['start']].strip()
                if before_table:
                    pre_chunks = self.standard_splitter.split_text(before_table)
                    chunks.extend(pre_chunks)
                
                # Handle the table as a special chunk
                table_content = table_section['content']
                
                if len(table_content) <= self.table_chunk_size:
                    # Table fits in one chunk, keep it together
                    chunks.append(table_content)
                else:
                    # Large table, split it intelligently
                    table_chunks = self._split_large_table(table_content)
                    chunks.extend(table_chunks)
                
                current_pos = table_section['end']
            
            # Process remaining text after last table
            remaining_text = text[current_pos:].strip()
            if remaining_text:
                remaining_chunks = self.standard_splitter.split_text(remaining_text)
                chunks.extend(remaining_chunks)
            
            return [chunk for chunk in chunks if chunk.strip()]
            
        except Exception as e:
            logger.error(f"Error in table-aware splitting: {e}")
            return self.standard_splitter.split_text(text)
    
    def _split_large_table(self, table_content: str) -> List[str]:
        """
        Split large tables while preserving structure
        
        Args:
            table_content: Table content that's too large for one chunk
            
        Returns:
            List of table chunks with preserved headers
        """
        try:
            lines = table_content.split('\n')
            
            # Find header information
            header_lines = []
            data_lines = []
            footer_lines = []
            
            in_data_section = False
            
            for line in lines:
                if line.startswith('===') and 'TABLE' in line and not line.startswith('=== END'):
                    header_lines.append(line)
                elif line.startswith('HEADERS:') or line == '---':
                    header_lines.append(line)
                elif line.startswith('=== END'):
                    footer_lines.append(line)
                    in_data_section = False
                elif line.startswith('ROW ') or in_data_section:
                    data_lines.append(line)
                    in_data_section = True
                else:
                    if not in_data_section:
                        header_lines.append(line)
                    else:
                        data_lines.append(line)
            
            # Create chunks with preserved headers
            header_text = '\n'.join(header_lines)
            footer_text = '\n'.join(footer_lines)
            
            chunks = []
            current_chunk_lines = [header_text] if header_text.strip() else []
            current_size = len(header_text)
            
            for line in data_lines:
                line_size = len(line) + 1  # +1 for newline
                
                if current_size + line_size > self.table_chunk_size and current_chunk_lines:
                    # Finalize current chunk
                    chunk_content = '\n'.join(current_chunk_lines)
                    if footer_text.strip():
                        chunk_content += '\n' + footer_text
                    chunks.append(chunk_content)
                    
                    # Start new chunk with header
                    current_chunk_lines = [header_text] if header_text.strip() else []
                    current_size = len(header_text)
                
                current_chunk_lines.append(line)
                current_size += line_size
            
            # Add final chunk
            if current_chunk_lines:
                chunk_content = '\n'.join(current_chunk_lines)
                if footer_text.strip():
                    chunk_content += '\n' + footer_text
                chunks.append(chunk_content)
            
            return chunks
            
        except Exception as e:
            logger.error(f"Error splitting large table: {e}")
            # Fallback to standard splitting
            return self.standard_splitter.split_text(table_content)



# Create global instance
document_service = DocumentService()
